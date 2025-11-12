from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_PATH = Path("docs/finance/MagicLawyer_Custos_Mensais.pptx")
BRAND_PRIMARY = RGBColor(25, 25, 112)  # deep indigo
BRAND_SECONDARY = RGBColor(0, 162, 174)  # teal accent
BRAND_TERTIARY = RGBColor(255, 140, 0)  # warm highlight
BRAND_BG = RGBColor(245, 247, 252)
TEXT_COLOR = RGBColor(38, 38, 38)

prs = Presentation()


def _apply_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_BG
    # Accent arc
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1.2),
        Inches(-1.2),
        Inches(4.5),
        Inches(4.5),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = BRAND_SECONDARY
    circle.line.color.rgb = BRAND_SECONDARY
    circle.shadow.inherit = False
    circle.rotation = 15


def _add_title(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = BRAND_PRIMARY
    if subtitle:
        sub = tf.add_paragraph()
        sub.text = subtitle
        sub.font.size = Pt(20)
        sub.font.color.rgb = RGBColor(70, 70, 70)


def _add_icon_card(slide, left, top, width, height, emoji, title, value, caption, accent=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = accent or BRAND_SECONDARY
    shape.shadow.inherit = True
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(26)
    p.font.color.rgb = accent or BRAND_SECONDARY
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRAND_PRIMARY
    p = tf.add_paragraph()
    p.text = value
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p = tf.add_paragraph()
    p.text = caption
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_COLOR


def _add_bullets(slide, left, top, width, height, items, size=18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.add_paragraph() if idx else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0


def _add_table(slide, left, top, width, height, data, header_fill=None, body_fill=None, col_widths=None):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for idx, col_w in enumerate(col_widths):
            table.columns[idx].width = col_w
    else:
        default_width = int(width) // cols
        for c in range(cols):
            table.columns[c].width = default_width
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12 if r else 13)
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = BRAND_PRIMARY
                else:
                    p.font.color.rgb = TEXT_COLOR
            cell.fill.solid()
            if r == 0 and header_fill:
                cell.fill.fore_color.rgb = header_fill
            elif r > 0 and body_fill:
                cell.fill.fore_color.rgb = body_fill if r % 2 else RGBColor(255, 255, 255)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)


def _add_section_label(slide, text):
    label = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(2.5), Inches(0.4))
    tf = label.text_frame
    tf.text = text.upper()
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BRAND_SECONDARY


# Slide 1 – capa
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "Magic Lawyer – 10/11/2025")
_add_title(slide, "Plano Financeiro Mensal", "História de investimento e assinatura para o parceiro certo")
_add_icon_card(
    slide,
    Inches(0.6),
    Inches(2.2),
    Inches(3.1),
    Inches(1.9),
    "💠",
    "Infra fixa",
    "R$ 1,678 mil",
    "Hosting, banco, realtime, DNS e contas Google",
)
_add_icon_card(
    slide,
    Inches(3.9),
    Inches(2.2),
    Inches(3.1),
    Inches(1.9),
    "📈",
    "Envelope variável",
    "R$ 2,35 mil",
    "Pagamentos, mensageria, storage extra, WhatsApp",
    accent=BRAND_PRIMARY,
)
_add_icon_card(
    slide,
    Inches(7.2),
    Inches(2.2),
    Inches(3.1),
    Inches(1.9),
    "🚀",
    "Break-even",
    "Mix 5-3-1",
    "5 Básico + 3 Pro + 1 Enterprise = R$ 4,6 mil",
    accent=BRAND_TERTIARY,
)

# Slide 2 – visão de investimento
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "snapshot")
_add_title(slide, "Visão Rápida", "Quanto precisamos por mês e como cada bloco se paga")
_add_bullets(
    slide,
    Inches(0.7),
    Inches(1.8),
    Inches(9.5),
    Inches(3.6),
    [
        "Capex recorrente estimado: **R$ 4,0 mil/mês** (infra obrigatória + variáveis moderadas).",
        "Foco em previsibilidade: contratos anuais com cloud + provisionamento de variáveis evita sustos.",
        "MRR alvo para break-even imediato: **R$ 1,7 mil** (7 Básico ou 3 Pro).",
        "Margem operacional desejada pós-break-even: **≥ R$ 2,5 mil/mês** para CAC e sucesso do cliente.",
        "Escala adicional vem de add-ons (WhatsApp oficial, storage premium, onboarding white-glove).",
    ],
)

# Slide 3 – Infraestrutura fixa
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "infraestrutura")
_add_title(slide, "Stack Obrigatório", "Serviços que mantêm o Magic Lawyer 24/7")
infra_table = [
    ["Serviço", "Função", "Plano", "R$/mês"],
    ["Vercel", "Hosting Next.js + crons", "Team Pro", "220"],
    ["Neon/Supabase", "PostgreSQL multi-tenant", "Pro 2 CU / 500 GB", "435"],
    ["Upstash Redis", "BullMQ, locks e cache", "Pro 100M cmds", "110"],
    ["Ably", "Realtime multi-tenant", "Business 3M msgs", "270"],
    ["Cloudinary", "Docs pesados e transformações", "Advanced 600 créditos", "545"],
    ["Google Workspace", "SMTP + contas core", "Business Starter (2)", "78"],
    ["Domínios + Cloudflare", "DNS, SSL e WAF", "Registros + CF Pro", "45"],
]
_add_table(slide, Inches(0.5), Inches(1.9), Inches(9.4), Inches(4), infra_table, header_fill=RGBColor(229, 235, 248), body_fill=RGBColor(241, 247, 255))

# Slide 4 – variáveis
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "variáveis")
_add_title(slide, "Envelope Variável", "Custos sensíveis a volume e engrenagens de receita")
var_table = [
    ["Item", "Métrica de cobrança", "Exemplo/mês", "Investimento (R$)"],
    ["Asaas", "Boleto/PIX/cartão", "200 boletos + 100 PIX + 80 cartões", "1.414"],
    ["Cloudinary extra", "Crédito adicional", "+300 créditos", "248"],
    ["Ably excedente", "Mensagens > 3M", "+5M msgs", "69"],
    ["Upstash overage", "Comandos adicionais", "+300k cmds", "3"],
    ["Backups S3/Wasabi", "200 GB", "Snapshots + assets", "25"],
    ["ngrok Pro", "QA de webhooks", "1 túnel", "88"],
    ["Resend fallback", "40k e-mails", "Starter + excedente", "154"],
    ["Meta Cloud API", "1.000 conversas", "Mix autenticação/utilidade", "350"],
]
_add_table(slide, Inches(0.5), Inches(1.9), Inches(9.4), Inches(4), var_table, header_fill=RGBColor(229, 248, 247), body_fill=RGBColor(236, 251, 250))

# Slide 5 – planos
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "receita")
_add_title(slide, "Oferta de Assinatura", "O que cada plano entrega e quanto cobramos")
plans_table = [
    ["Plano", "Perfil", "Limites inclusos", "Preço mensal", "Preço anual"],
    ["Básico", "Até 3 usuários", "50 processos / 1 GB / 500 docs", "R$ 249", "R$ 2.490"],
    ["Pro", "Até 10 usuários", "200 processos / 5 GB / 2k docs", "R$ 699", "R$ 6.990"],
    ["Enterprise", "Até 50 usuários", "1.000 processos / 20 GB / 10k docs", "R$ 1.299", "R$ 12.990"],
    ["Ultra", "Sob demanda", "Limites customizados + gerente", "R$ 2.490+", "Sob consulta"],
]
_add_table(slide, Inches(0.4), Inches(1.9), Inches(9.6), Inches(3.9), plans_table, header_fill=RGBColor(253, 237, 218), body_fill=RGBColor(255, 248, 237))
_add_bullets(
    slide,
    Inches(0.5),
    Inches(5.1),
    Inches(9.4),
    Inches(1.0),
    [
        "Add-ons: WhatsApp oficial, storage adicional, blocos de documentos e onboarding premium.",
        "Upsell planejado após adoção do Pro (integrações PJe, API e automações).",
    ],
    size=14,
)

# Slide 6 – break-even
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "rentabilidade")
_add_title(slide, "Break-even & Payback", "Como a assinatura cobre o investimento mensal")
_add_icon_card(
    slide,
    Inches(0.5),
    Inches(1.8),
    Inches(3.3),
    Inches(1.9),
    "🎯",
    "Meta 5-3-1",
    "R$ 4,641",
    "5 Básico + 3 Pro + 1 Enterprise pagam o mês e sobram R$ 2,963",
)
_add_icon_card(
    slide,
    Inches(3.9),
    Inches(1.8),
    Inches(3.3),
    Inches(1.9),
    "⚡️",
    "Payback imediato",
    "R$ 1,7 mil",
    "7 Básico ou 3 Pro já cobrem a infra fixa de R$ 1,678",
    accent=BRAND_PRIMARY,
)
_add_icon_card(
    slide,
    Inches(7.3),
    Inches(1.8),
    Inches(3.3),
    Inches(1.9),
    "💸",
    "Margem alvo",
    "≥ R$ 2,5 mil",
    "Reserva mensal para CAC, marketing e sucesso do cliente",
    accent=BRAND_TERTIARY,
)
_add_table(
    slide,
    Inches(0.5),
    Inches(3.9),
    Inches(9.6),
    Inches(2.1),
    [
        ["Mix", "Receita", "Margem sobre fixo", "Comentário"],
        ["5B + 4P + 1E", "R$ 5.340", "R$ 3.662", "Foco em Pro acelera margem"],
        ["5B + 5P", "R$ 4.740", "R$ 3.062", "Sem Enterprise ainda cobre folgado"],
        ["5B + 2P + 2E", "R$ 5.241", "R$ 3.563", "Enterprise libera caixa para marketing"],
    ],
    header_fill=RGBColor(228, 236, 255),
    body_fill=RGBColor(239, 244, 255),
)

# Slide 7 – escala
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "escala")
_add_title(slide, "Projeção por Estágio", "Infra + variáveis versus crescimento de tenants")
scale_table = [
    ["Estágio", "Tenants", "Usuários", "Processos/mês", "Receita EUA (Asaas)", "Opex variável", "Custo total"],
    ["Lançamento", "3", "60", "150", "R$ 90 mil", "R$ 450", "R$ 2.130"],
    ["Crescimento", "10", "250", "600", "R$ 360 mil", "R$ 1.650", "R$ 3.330"],
    ["Escala regional", "25", "700", "1.800", "R$ 1,1 milhão", "R$ 4.900", "R$ 6.580"],
]
_add_table(slide, Inches(0.4), Inches(1.9), Inches(9.6), Inches(3.8), scale_table, header_fill=RGBColor(227, 245, 255), body_fill=RGBColor(237, 250, 255))
_add_bullets(
    slide,
    Inches(0.5),
    Inches(5.0),
    Inches(9.4),
    Inches(1.1),
    [
        "Receita Asaas assume ticket de R$ 600 e taxa média de 3%.",
        "Acima de 25 tenants: planejar Postgres dedicado e Redis dimensionado.",
    ],
    size=14,
)

# Slide 8 – próximos passos
slide = prs.slides.add_slide(prs.slide_layouts[6])
_apply_bg(slide)
_add_section_label(slide, "ação")
_add_title(slide, "Próximos Passos", "Onde investir energia após o cheque")
_add_bullets(
    slide,
    Inches(0.6),
    Inches(1.9),
    Inches(4.6),
    Inches(3.5),
    [
        "Rever câmbio e contratos com cloud trimestralmente.",
        "Rodar dashboards públicos do “5-3-1” para toda a operação.",
        "Pacotar add-ons premium (WhatsApp oficial, storage extra, onboarding white-glove).",
    ],
)
_add_bullets(
    slide,
    Inches(5.3),
    Inches(1.9),
    Inches(4.6),
    Inches(3.5),
    [
        "Gatilho de reinvestimento: margem ≥ R$ 2,5 mil → CAC e conteúdo.",
        "Atualizar documento financeiro sempre que entrar novo serviço pago.",
        "Preparar métricas de adoção para próxima rodada (churn, NRR, payback CAC).",
    ],
)

prs.save(OUTPUT_PATH)
print(f"Apresentação gerada em {OUTPUT_PATH}")
